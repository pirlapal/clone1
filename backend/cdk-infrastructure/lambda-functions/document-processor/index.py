import json
import boto3
import os
import tempfile
import logging
from typing import Dict, Any, Optional
from io import BytesIO
import uuid

# Document processing libraries
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.utils import ImageReader
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from PIL import Image as PILImage
import PyPDF2

# Configure logging
logger = logging.getLogger()
logger.setLevel(logging.INFO)

# Initialize AWS clients
s3_client = boto3.client('s3')
bedrock_agent = boto3.client('bedrock-agent')

def handler(event: Dict[str, Any], context) -> Dict[str, Any]:
    """
    Process uploaded documents (PPT to PDF conversion, PDF optimization)
    Triggered by S3 uploads
    """
    logger.info(f"Processing document event: {json.dumps(event, default=str)}")
    
    try:
        processed_files = []
        
        # Extract S3 event details
        for record in event['Records']:
            bucket = record['s3']['bucket']['name']
            key = record['s3']['object']['key']
            
            logger.info(f"Processing file: {key} from bucket: {bucket}")
            
            # Skip if already in processed folder
            if key.startswith('processed/'):
                logger.info(f"File {key} already in processed folder, skipping")
                continue
            
            # Process based on file type
            processed_key = None
            
            if key.lower().endswith(('.ppt', '.pptx')):
                processed_key = process_powerpoint(bucket, key)
            elif key.lower().endswith('.pdf'):
                processed_key = process_pdf(bucket, key)
            elif key.lower().endswith(('.doc', '.docx')):
                processed_key = process_word_document(bucket, key)
            else:
                logger.warning(f"Unsupported file type: {key}")
                continue
            
            if processed_key:
                processed_files.append(processed_key)
                
                # Trigger knowledge base sync
                trigger_knowledge_base_sync()
        
        return {
            'statusCode': 200,
            'body': json.dumps({
                'message': 'Documents processed successfully',
                'processed_files': processed_files
            })
        }
        
    except Exception as e:
        logger.error(f"Error processing documents: {str(e)}", exc_info=True)
        return {
            'statusCode': 500,
            'body': json.dumps({'error': str(e)})
        }

def process_powerpoint(bucket: str, key: str) -> Optional[str]:
    """Convert PowerPoint to PDF optimized for Bedrock Data Automation"""
    logger.info(f"Converting PowerPoint to PDF with Data Automation optimization: {key}")
    
    try:
        # Download PPT file
        with tempfile.NamedTemporaryFile(suffix='.pptx') as temp_ppt:
            s3_client.download_fileobj(bucket, key, temp_ppt)
            temp_ppt.flush()
            
            # Load presentation
            prs = Presentation(temp_ppt.name)
            
            # Create PDF optimized for Bedrock Data Automation
            with tempfile.NamedTemporaryFile(suffix='.pdf') as temp_pdf:
                doc = SimpleDocTemplate(temp_pdf.name, pagesize=A4)
                styles = getSampleStyleSheet()
                story = []
                
                # Document metadata for Data Automation
                filename = os.path.basename(key)
                title_style = ParagraphStyle(
                    'DocumentTitle',
                    parent=styles['Title'],
                    fontSize=24,
                    spaceAfter=30,
                    alignment=1,
                    textColor='black'
                )
                
                # Add document header with metadata
                story.append(Paragraph(f"DOCUMENT: {filename}", title_style))
                story.append(Paragraph(f"SOURCE: PowerPoint Presentation", styles['Normal']))
                story.append(Paragraph(f"PROCESSED: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", styles['Normal']))
                story.append(Paragraph(f"TOTAL_SLIDES: {len(prs.slides)}", styles['Normal']))
                story.append(Spacer(1, 30))
                
                # Process each slide with enhanced structure for Data Automation
                for i, slide in enumerate(prs.slides):
                    # Slide header
                    slide_header_style = ParagraphStyle(
                        'SlideHeader',
                        parent=styles['Heading1'],
                        fontSize=18,
                        spaceAfter=15,
                        textColor='darkblue',
                        borderWidth=1,
                        borderColor='lightgray',
                        borderPadding=5
                    )
                    story.append(Paragraph(f"SLIDE {i + 1} OF {len(prs.slides)}", slide_header_style))
                    
                    # Extract and structure slide content
                    slide_content = extract_structured_slide_content(slide)
                    
                    # Add slide title if found
                    if slide_content.get('title'):
                        title_style = ParagraphStyle(
                            'SlideTitle',
                            parent=styles['Heading2'],
                            fontSize=16,
                            spaceAfter=10,
                            textColor='darkgreen'
                        )
                        story.append(Paragraph(f"TITLE: {slide_content['title']}", title_style))
                    
                    # Add main content
                    if slide_content.get('content'):
                        content_style = ParagraphStyle(
                            'SlideContent',
                            parent=styles['Normal'],
                            fontSize=12,
                            spaceAfter=8,
                            leftIndent=20
                        )
                        for content_item in slide_content['content']:
                            story.append(Paragraph(f"• {content_item}", content_style))
                    
                    # Add tables if present
                    if slide_content.get('tables'):
                        table_header_style = ParagraphStyle(
                            'TableHeader',
                            parent=styles['Heading3'],
                            fontSize=14,
                            spaceAfter=5,
                            textColor='purple'
                        )
                        story.append(Paragraph("TABLES:", table_header_style))
                        
                        for table_idx, table_data in enumerate(slide_content['tables']):
                            story.append(Paragraph(f"Table {table_idx + 1}:", styles['Normal']))
                            for row in table_data:
                                story.append(Paragraph(f"  {' | '.join(row)}", styles['Normal']))
                            story.append(Spacer(1, 10))
                    
                    # Add speaker notes if present
                    if slide_content.get('notes'):
                        notes_style = ParagraphStyle(
                            'SpeakerNotes',
                            parent=styles['Normal'],
                            fontSize=10,
                            textColor='gray',
                            leftIndent=20,
                            spaceAfter=10
                        )
                        story.append(Paragraph("SPEAKER NOTES:", styles['Heading4']))
                        story.append(Paragraph(slide_content['notes'], notes_style))
                    
                    # Add slide separator
                    story.append(Spacer(1, 20))
                    story.append(Paragraph("─" * 80, styles['Normal']))
                    story.append(Spacer(1, 20))
                
                # Build PDF
                doc.build(story)
                temp_pdf.flush()
                
                # Upload to processed folder with enhanced metadata
                processed_key = f"processed/{os.path.splitext(os.path.basename(key))[0]}.pdf"
                
                with open(temp_pdf.name, 'rb') as pdf_file:
                    s3_client.upload_fileobj(
                        pdf_file,
                        bucket,
                        processed_key,
                        ExtraArgs={
                            'ContentType': 'application/pdf',
                            'Metadata': {
                                'original-file': key,
                                'conversion-type': 'ppt-to-pdf-data-automation',
                                'processed-timestamp': str(int(time.time())),
                                'slide-count': str(len(prs.slides)),
                                'document-type': 'presentation',
                                'parsing-strategy': 'bedrock-data-automation',
                                'content-structure': 'hierarchical'
                            }
                        }
                    )
                
                logger.info(f"Successfully converted {key} to {processed_key} with Data Automation optimization")
                return processed_key
                
    except Exception as e:
        logger.error(f"Error converting PowerPoint {key}: {str(e)}", exc_info=True)
        return None

def extract_structured_slide_content(slide) -> dict:
    """Extract structured content from a PowerPoint slide for Data Automation"""
    content = {
        'title': None,
        'content': [],
        'tables': [],
        'notes': None
    }
    
    try:
        # Extract speaker notes
        if hasattr(slide, 'notes_slide') and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                content['notes'] = notes_text
        
        # Process shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                
                # Identify title (usually the first large text or specific placeholder)
                if not content['title'] and (
                    hasattr(shape, 'placeholder_format') and 
                    shape.placeholder_format.type == 1  # Title placeholder
                ):
                    content['title'] = text
                elif not content['title'] and len(text) < 100:  # Assume short text is title
                    content['title'] = text
                else:
                    # Split multi-line content
                    lines = [line.strip() for line in text.split('\n') if line.strip()]
                    content['content'].extend(lines)
            
            elif shape.has_table:
                # Extract table content with structure
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = cell.text.strip() if cell.text else ""
                        row_data.append(cell_text)
                    if any(row_data):  # Only add non-empty rows
                        table_data.append(row_data)
                
                if table_data:
                    content['tables'].append(table_data)
    
    except Exception as e:
        logger.warning(f"Error extracting structured slide content: {str(e)}")
    
    return content

def extract_text_from_slide(slide) -> list:
    """Extract text content from a PowerPoint slide"""
    texts = []
    
    try:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            elif shape.has_table:
                # Extract table content
                table_text = extract_table_text(shape.table)
                if table_text:
                    texts.extend(table_text)
    except Exception as e:
        logger.warning(f"Error extracting text from slide: {str(e)}")
    
    return texts

def extract_table_text(table) -> list:
    """Extract text from PowerPoint table"""
    texts = []
    
    try:
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                if cell.text.strip():
                    row_text.append(cell.text.strip())
            if row_text:
                texts.append(" | ".join(row_text))
    except Exception as e:
        logger.warning(f"Error extracting table text: {str(e)}")
    
    return texts

def process_pdf(bucket: str, key: str) -> Optional[str]:
    """Process and optimize PDF files for Bedrock Data Automation"""
    logger.info(f"Processing PDF for Data Automation: {key}")
    
    try:
        # Download PDF file
        with tempfile.NamedTemporaryFile(suffix='.pdf') as temp_input:
            s3_client.download_fileobj(bucket, key, temp_input)
            temp_input.flush()
            
            # Validate and optimize PDF for Data Automation
            with open(temp_input.name, 'rb') as pdf_file:
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                
                # Check if PDF is readable
                if len(pdf_reader.pages) == 0:
                    logger.error(f"PDF {key} has no readable pages")
                    return None
                
                # Extract metadata for Data Automation
                metadata = pdf_reader.metadata or {}
                
                # Create optimized PDF with enhanced metadata
                with tempfile.NamedTemporaryFile(suffix='.pdf') as temp_output:
                    pdf_writer = PyPDF2.PdfWriter()
                    
                    # Add metadata for Data Automation
                    pdf_writer.add_metadata({
                        '/Title': metadata.get('/Title', os.path.basename(key)),
                        '/Author': metadata.get('/Author', 'Unknown'),
                        '/Subject': metadata.get('/Subject', 'Document processed for Bedrock Data Automation'),
                        '/Creator': 'iECHO Document Processor',
                        '/Producer': 'Bedrock Data Automation Pipeline',
                        '/Keywords': f"bedrock,data-automation,{os.path.splitext(os.path.basename(key))[0]}"
                    })
                    
                    # Copy pages and extract text for validation
                    total_text_length = 0
                    for page_num, page in enumerate(pdf_reader.pages):
                        try:
                            pdf_writer.add_page(page)
                            
                            # Extract text to validate content
                            page_text = page.extract_text()
                            total_text_length += len(page_text) if page_text else 0
                            
                        except Exception as e:
                            logger.warning(f"Error processing page {page_num} in {key}: {str(e)}")
                    
                    # Write optimized PDF
                    with open(temp_output.name, 'wb') as output_file:
                        pdf_writer.write(output_file)
                    
                    temp_output.flush()
                    
                    # Upload to processed folder with Data Automation metadata
                    processed_key = f"processed/{os.path.basename(key)}"
                    
                    with open(temp_output.name, 'rb') as optimized_pdf:
                        s3_client.upload_fileobj(
                            optimized_pdf,
                            bucket,
                            processed_key,
                            ExtraArgs={
                                'ContentType': 'application/pdf',
                                'Metadata': {
                                    'original-file': key,
                                    'conversion-type': 'pdf-optimization-data-automation',
                                    'processed-timestamp': str(int(time.time())),
                                    'page-count': str(len(pdf_reader.pages)),
                                    'document-type': 'pdf',
                                    'parsing-strategy': 'bedrock-data-automation',
                                    'content-length': str(total_text_length),
                                    'title': str(metadata.get('/Title', 'Unknown')),
                                    'author': str(metadata.get('/Author', 'Unknown'))
                                }
                            }
                        )
                    
                    logger.info(f"Successfully processed PDF {key} to {processed_key} for Data Automation")
                    return processed_key
                    
    except Exception as e:
        logger.error(f"Error processing PDF {key}: {str(e)}", exc_info=True)
        return None

def process_word_document(bucket: str, key: str) -> Optional[str]:
    """Convert Word documents to PDF"""
    logger.info(f"Converting Word document to PDF: {key}")
    
    try:
        from docx import Document
        
        # Download Word file
        with tempfile.NamedTemporaryFile(suffix='.docx') as temp_docx:
            s3_client.download_fileobj(bucket, key, temp_docx)
            temp_docx.flush()
            
            # Load document
            doc = Document(temp_docx.name)
            
            # Create PDF
            with tempfile.NamedTemporaryFile(suffix='.pdf') as temp_pdf:
                pdf_doc = SimpleDocTemplate(temp_pdf.name, pagesize=A4)
                styles = getSampleStyleSheet()
                story = []
                
                # Title
                filename = os.path.basename(key)
                title_style = ParagraphStyle(
                    'CustomTitle',
                    parent=styles['Heading1'],
                    fontSize=20,
                    spaceAfter=20,
                    alignment=1
                )
                story.append(Paragraph(f"Document: {filename}", title_style))
                story.append(Spacer(1, 20))
                
                # Process paragraphs
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        # Determine style based on paragraph formatting
                        if paragraph.style.name.startswith('Heading'):
                            story.append(Paragraph(paragraph.text, styles['Heading2']))
                        else:
                            story.append(Paragraph(paragraph.text, styles['Normal']))
                        story.append(Spacer(1, 6))
                
                # Process tables
                for table in doc.tables:
                    story.append(Paragraph("Table:", styles['Heading3']))
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            story.append(Paragraph(" | ".join(row_text), styles['Normal']))
                    story.append(Spacer(1, 12))
                
                # Build PDF
                pdf_doc.build(story)
                temp_pdf.flush()
                
                # Upload to processed folder
                processed_key = f"processed/{os.path.splitext(os.path.basename(key))[0]}.pdf"
                
                with open(temp_pdf.name, 'rb') as pdf_file:
                    s3_client.upload_fileobj(
                        pdf_file,
                        bucket,
                        processed_key,
                        ExtraArgs={
                            'ContentType': 'application/pdf',
                            'Metadata': {
                                'original-file': key,
                                'conversion-type': 'docx-to-pdf',
                                'processed-timestamp': str(int(time.time()))
                            }
                        }
                    )
                
                logger.info(f"Successfully converted {key} to {processed_key}")
                return processed_key
                
    except Exception as e:
        logger.error(f"Error converting Word document {key}: {str(e)}", exc_info=True)
        return None

def trigger_knowledge_base_sync():
    """Trigger Bedrock Knowledge Base synchronization"""
    try:
        kb_id = os.environ.get('KNOWLEDGE_BASE_ID')
        ds_id = os.environ.get('DATA_SOURCE_ID')
        
        if not kb_id or not ds_id:
            logger.warning("Knowledge Base ID or Data Source ID not found in environment")
            return
        
        response = bedrock_agent.start_ingestion_job(
            knowledgeBaseId=kb_id,
            dataSourceId=ds_id,
            description=f"Auto-sync triggered by document processing at {int(time.time())}"
        )
        
        job_id = response['ingestionJob']['ingestionJobId']
        logger.info(f"Started knowledge base ingestion job: {job_id}")
        
    except Exception as e:
        logger.error(f"Error triggering knowledge base sync: {str(e)}", exc_info=True)

# Import time module
import time
